"""
Generate Agentic AI Testing Architecture PowerPoint Presentation
17 Slides — Professional, clean, structured for presenting
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BG     = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BLUE  = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_TEAL = RGBColor(0x00, 0xC9, 0xA7)
LIGHT_GRAY  = RGBColor(0xF0, 0xF0, 0xF5)
MID_GRAY    = RGBColor(0x6B, 0x7B, 0x8D)
DARK_GRAY   = RGBColor(0x2D, 0x3A, 0x4A)
RED_ALERT   = RGBColor(0xE7, 0x4C, 0x3C)
GREEN_OK    = RGBColor(0x2E, 0xCC, 0x71)
ORANGE_WARN = RGBColor(0xF3, 0x9C, 0x12)
PURPLE      = RGBColor(0x8E, 0x44, 0xAD)
SECTION_BG  = RGBColor(0x0D, 0x15, 0x2D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helper Functions ───────────────────────────────────────────

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        try:
            from lxml import etree
            from pptx.oxml.ns import qn
            sp_elem = shape._element
            solid_fill = sp_elem.find('.//' + qn('a:solidFill'))
            if solid_fill is not None:
                srgb = solid_fill.find(qn('a:srgbClr'))
                if srgb is not None:
                    a_elem = etree.SubElement(srgb, qn('a:alpha'))
                    a_elem.set('val', str(int(alpha * 1000)))
        except Exception:
            pass  # Skip alpha if XML manipulation fails
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT_TEAL, font_name="Calibri",
                    spacing=6, bold_prefix=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle bold prefix (text before " — " or " - " at start)
        if bold_prefix and (" — " in item or " – " in item):
            sep = " — " if " — " in item else " – "
            parts = item.split(sep, 1)
            run1 = p.add_run()
            run1.text = "  " + parts[0] + sep
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.bold = True
            run1.font.name = font_name
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.bold = False
            run2.font.name = font_name
        else:
            run = p.add_run()
            run.text = "  " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = False
            run.font.name = font_name

        # Add bullet character
        bullet_run = p.add_run()
        bullet_run.text = ""
        p.space_after = Pt(spacing)
        p.space_before = Pt(spacing)

        # Bullet via XML
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        buNone = pPr.findall(qn('a:buNone'))
        for bn in buNone:
            pPr.remove(bn)
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '▸'})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = buClr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (bullet_color[0], bullet_color[1], bullet_color[2]) if isinstance(bullet_color, tuple) else str(bullet_color).replace('#','')})
        buClr.append(srgb)
        pPr.append(buClr)
        pPr.append(buChar)
        buSz = pPr.makeelement(qn('a:buSzPct'), {'val': '100000'})
        pPr.append(buSz)

    return txBox


def add_table(slide, left, top, width, height, headers, rows, 
              header_color=ACCENT_BLUE, font_size=13):
    cols = len(headers)
    row_count = len(rows) + 1
    table_shape = slide.shapes.add_table(row_count, cols, left, top, width, height)
    table = table_shape.table

    col_width = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_width

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1E, 0x2D, 0x4A)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size - 1)
            p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_accent_bar(slide, left, top, width=Inches(0.06), height=Inches(1.0), color=ACCENT_TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, items, accent=ACCENT_TEAL,
             font_size=13, title_size=15):
    # Card background
    card = add_shape_bg(slide, left, top, width, height, RGBColor(0x1E, 0x2D, 0x4A))
    # Accent top bar
    add_shape_bg(slide, left, top, width, Inches(0.05), accent)
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 title, font_size=title_size, color=accent, bold=True)
    # Items
    y = top + Inches(0.55)
    for item in items:
        add_text_box(slide, left + Inches(0.2), y, width - Inches(0.4), Inches(0.3),
                     "▸ " + item, font_size=font_size, color=RGBColor(0xCC, 0xCC, 0xCC))
        y += Inches(0.28)
    return card


def add_slide_number(slide, num, total=17):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


def add_footer_line(slide):
    add_shape_bg(slide, Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.01), ACCENT_BLUE)


def slide_header(slide, title, subtitle=None, num=1):
    """Standard slide header with accent bar"""
    add_accent_bar(slide, Inches(0.5), Inches(0.4), Inches(0.06), Inches(0.6), ACCENT_TEAL)
    add_text_box(slide, Inches(0.75), Inches(0.3), Inches(10), Inches(0.6),
                 title, font_size=30, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.75), Inches(0.85), Inches(10), Inches(0.4),
                     subtitle, font_size=16, color=MID_GRAY)
    add_slide_number(slide, num)
    add_footer_line(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, SECTION_BG)

# Gradient overlay
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RGBColor(0x00, 0x2B, 0x5C), alpha=40)

# Top accent line
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_TEAL)

# Title block
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(1.2),
             "Agentic AI Testing Architecture",
             font_size=44, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(2.5), Inches(11), Inches(0.8),
             "for Automated Tool Validation",
             font_size=32, color=LIGHT_BLUE)

# Divider
add_shape_bg(slide, Inches(1.0), Inches(3.5), Inches(3), Inches(0.04), ACCENT_TEAL)

# Subtitle
add_text_box(slide, Inches(1.0), Inches(3.9), Inches(11), Inches(0.6),
             "A Meta-Testing Platform — It tests applications, and it tests itself.",
             font_size=20, color=RGBColor(0xAA, 0xBB, 0xCC))

# Three pillars
pillar_data = [
    ("ACCURACY", "AI quality is measured,\nnot assumed", ACCENT_TEAL),
    ("AUTONOMY", "Minimal human intervention\nfor routine testing", LIGHT_BLUE),
    ("TRUST", "Every AI decision is\nauditable & traceable", ACCENT_BLUE),
]
x_start = Inches(1.0)
for i, (title, desc, color) in enumerate(pillar_data):
    x = x_start + Inches(i * 3.8)
    add_shape_bg(slide, x, Inches(4.8), Inches(3.3), Inches(1.6), RGBColor(0x14, 0x1E, 0x3A))
    add_shape_bg(slide, x, Inches(4.8), Inches(3.3), Inches(0.05), color)
    add_text_box(slide, x + Inches(0.2), Inches(5.0), Inches(2.9), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(5.4), Inches(2.9), Inches(0.8),
                 desc, font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

add_text_box(slide, Inches(1.0), Inches(6.7), Inches(6), Inches(0.4),
             "Testing Architect Presentation  |  February 2026",
             font_size=12, color=MID_GRAY)
add_slide_number(slide, 1)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Problem Statement", "Why We Need an Agentic AI Testing Platform", 2)

# Four problem cards
problems = [
    ("Manual Test Design\nDoesn't Scale", "QA becomes bottleneck as\nsprints accelerate", "📈", RED_ALERT),
    ("Jira Stories Are\nAmbiguous", "35% missing AC\n60% miss negative scenarios", "📝", ORANGE_WARN),
    ("Automation Scripts\nBreak Frequently", "42% failures from\nbroken locators alone", "🔧", RED_ALERT),
    ("No Confidence in\nAI-Generated Tests", "No validation framework\nfor hallucination or coverage", "🤖", ORANGE_WARN),
]

for i, (title, desc, icon, color) in enumerate(problems):
    x = Inches(0.5) + Inches(i * 3.15)
    y = Inches(1.6)
    add_shape_bg(slide, x, y, Inches(2.95), Inches(2.3), RGBColor(0x1E, 0x2D, 0x4A))
    add_shape_bg(slide, x, y, Inches(2.95), Inches(0.05), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.6), Inches(0.8),
                 title, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.1), Inches(2.6), Inches(0.9),
                 desc, font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

# Goal section
add_shape_bg(slide, Inches(0.5), Inches(4.3), Inches(12.333), Inches(2.3), RGBColor(0x0D, 0x19, 0x33))
add_accent_bar(slide, Inches(0.5), Inches(4.3), Inches(0.06), Inches(2.3), GREEN_OK)

add_text_box(slide, Inches(0.9), Inches(4.5), Inches(4), Inches(0.4),
             "THE GOAL", font_size=20, color=GREEN_OK, bold=True)
add_text_box(slide, Inches(0.9), Inches(4.95), Inches(11.5), Inches(0.5),
             "Build a self-learning testing platform that:", font_size=17, color=WHITE)

goals = [
    "Scales test design without scaling the team",
    "Detects and flags ambiguous requirements automatically",
    "Self-heals broken automation scripts",
    "Measures AI quality with empirical metrics, not assumptions",
]
y = Inches(5.5)
for g in goals:
    add_text_box(slide, Inches(1.2), y, Inches(11), Inches(0.32),
                 "✓  " + g, font_size=15, color=RGBColor(0xCC, 0xDD, 0xCC))
    y += Inches(0.32)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — HIGH-LEVEL ARCHITECTURE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "High-Level Architecture", "5 Independent, Testable Layers", 3)

layers = [
    ("LAYER 5", "Control Plane", "Supervisor Agent — Coordinates, decides, retries, escalates", PURPLE),
    ("LAYER 4", "Analysis & Metrics", "AI Metrics + Execution Metrics — Measures everything", ACCENT_BLUE),
    ("LAYER 3", "Automation & Execution", "Script Agent → Execution Engine → RCA Agent", LIGHT_BLUE),
    ("LAYER 2", "Agentic AI", "Requirement Agent → Test Case Agent → Feedback Loop", ACCENT_TEAL),
    ("LAYER 1", "Input Layer", "Jira Connector → Parser → Validator → Normalizer", GREEN_OK),
]

for i, (layer_num, name, desc, color) in enumerate(layers):
    y = Inches(1.5) + Inches(i * 1.05)
    # Layer background
    add_shape_bg(slide, Inches(0.5), y, Inches(12.333), Inches(0.9), RGBColor(0x1A, 0x2A, 0x44))
    # Color accent on left
    add_shape_bg(slide, Inches(0.5), y, Inches(0.08), Inches(0.9), color)
    # Layer number badge
    add_shape_bg(slide, Inches(0.8), y + Inches(0.15), Inches(1.2), Inches(0.6), color)
    add_text_box(slide, Inches(0.8), y + Inches(0.2), Inches(1.2), Inches(0.5),
                 layer_num, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Name
    add_text_box(slide, Inches(2.2), y + Inches(0.1), Inches(3), Inches(0.4),
                 name, font_size=18, color=WHITE, bold=True)
    # Description
    add_text_box(slide, Inches(2.2), y + Inches(0.48), Inches(10), Inches(0.4),
                 desc, font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

# Key insight
add_shape_bg(slide, Inches(0.5), Inches(6.85), Inches(12.333), Inches(0.01), ACCENT_BLUE)
add_text_box(slide, Inches(0.5), Inches(6.45), Inches(12), Inches(0.4),
             "Architect's View:  Each layer is independently testable. Data flows down, feedback flows up.",
             font_size=14, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — AGENTIC AI DESIGN PHILOSOPHY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Agentic AI Design Philosophy", '"Multiple specialized agents coordinated by a supervisor agent"', 4)

# Three pillars
pillars = [
    ("Autonomous Execution", ACCENT_TEAL,
     ["Agents act without step-by-step instructions",
      "Goal-driven, not rule-driven",
      "Pipeline runs end-to-end autonomously"]),
    ("Decision-Making Capability", LIGHT_BLUE,
     ["Each agent decides and logs its choices",
      "Confidence scores enable routing",
      "Decisions are explainable & reversible"]),
    ("Feedback-Driven Improvement", ACCENT_BLUE,
     ["Execution results tune LLM prompts",
      "Thresholds adjust based on outcomes",
      "System improves with every cycle"]),
]

for i, (title, color, items) in enumerate(pillars):
    x = Inches(0.5) + Inches(i * 4.15)
    y = Inches(1.6)
    add_shape_bg(slide, x, y, Inches(3.9), Inches(3.2), RGBColor(0x1A, 0x2A, 0x44))
    add_shape_bg(slide, x, y, Inches(3.9), Inches(0.05), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(0.5),
                 title, font_size=18, color=color, bold=True)
    iy = y + Inches(0.8)
    for item in items:
        add_text_box(slide, x + Inches(0.2), iy, Inches(3.5), Inches(0.35),
                     "▸  " + item, font_size=13, color=RGBColor(0xCC, 0xCC, 0xCC))
        iy += Inches(0.38)

# Comparison box
add_shape_bg(slide, Inches(0.5), Inches(5.1), Inches(12.333), Inches(1.6), RGBColor(0x0D, 0x19, 0x33))
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(5), Inches(0.4),
             "WHY MULTI-AGENT (NOT MONOLITHIC)?", font_size=16, color=ORANGE_WARN, bold=True)

reasons = [
    "Each agent fails independently — no single point of failure",
    "Each agent testable in isolation with its own golden dataset",
    "Agents scale independently (3 execution agents, 1 RCA agent)",
    "Clear responsibility boundaries make debugging tractable",
]
y = Inches(5.65)
for r in reasons:
    add_text_box(slide, Inches(1.0), y, Inches(11), Inches(0.3),
                 "✓  " + r, font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC))
    y += Inches(0.3)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — MODULE 1: JIRA INGESTION
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Module 1: Jira Ingestion & Validation", "The Gateway — Where Everything Begins", 5)

# What it does
add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.2),
         "WHAT IT DOES", [
             "Reads Jira story by ID or batch (OAuth/API Token)",
             "Validates structure — required fields, formats",
             "Normalizes content — strip HTML, resolve macros",
             "Outputs canonical JSON for downstream agents",
         ], accent=ACCENT_TEAL)

# Testing focus
add_card(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.2),
         "TESTING FOCUS", [
             "Empty/malformed stories → Flag and block",
             "Expired tokens / no access → Clear auth errors",
             "Rate limits (429) → Backoff and retry",
             "XSS payloads / special chars → Sanitized",
         ], accent=RED_ALERT)

# Metrics table
add_table(slide, Inches(0.5), Inches(4.1), Inches(12.333), Inches(1.6),
          ["Metric", "Target", "Alert Threshold"],
          [
              ["Ingestion Success Rate", ">= 98%", "< 95%"],
              ["Parsing Error Rate", "< 2%", "> 5%"],
              ["Validation Pass Rate", ">= 90%", "< 85%"],
              ["Avg Ingestion Latency", "< 2 seconds", "> 5 seconds"],
          ])

add_text_box(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.4),
             "Architect's Rule:  If garbage enters here, every downstream agent produces garbage.",
             font_size=14, color=ORANGE_WARN, bold=False)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — MODULE 2: REQUIREMENT UNDERSTANDING
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Module 2: Requirement Understanding Agent", "Turning Stories into Testable Knowledge", 6)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.9), Inches(2.8),
         "RESPONSIBILITIES", [
             "Extract acceptance criteria",
             "  (Gherkin, bullets, prose)",
             "Identify implicit business rules",
             "Detect ambiguity & flag vague",
             "  requirements for review",
         ], accent=ACCENT_TEAL)

add_card(slide, Inches(4.7), Inches(1.5), Inches(3.9), Inches(2.8),
         "TESTING STRATEGY", [
             "Golden Story Comparison",
             "  — compare vs human-verified",
             "Hallucination Detection",
             "  — no invented criteria",
             "Ambiguity Flag Accuracy",
         ], accent=LIGHT_BLUE)

add_card(slide, Inches(8.9), Inches(1.5), Inches(3.9), Inches(2.8),
         "AMBIGUITY EXAMPLES", [
             '✗ "should work correctly"',
             '✗ "handle errors appropriately"',
             '✗ "response time acceptable"',
             '✓ "login with email & password"',
             '✓ "display name on dashboard"',
         ], accent=ORANGE_WARN)

# Key metric
add_shape_bg(slide, Inches(0.5), Inches(4.7), Inches(12.333), Inches(1.8), RGBColor(0x0D, 0x19, 0x33))
add_text_box(slide, Inches(0.8), Inches(4.85), Inches(4), Inches(0.4),
             "KEY METRICS", font_size=18, color=ACCENT_TEAL, bold=True)

add_table(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.0),
          ["Metric", "Target", "Alert"],
          [
              ["Requirement Interpretation Accuracy", ">= 85%", "< 75%"],
              ["Hallucination Rate", "< 5%", "> 8%"],
          ])


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — MODULE 3: TEST CASE DESIGN
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Module 3: Test Case Design Agent", "From Requirements to Comprehensive Test Cases", 7)

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.0),
         "WHAT IT GENERATES", [
             "Positive path tests — happy flow for each AC",
             "Negative path tests — invalid inputs, unauthorized access",
             "Edge cases — boundary values, empty inputs, max lengths",
             "Risk-based priority (P0-P3) per test case",
         ], accent=ACCENT_TEAL)

add_card(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.0),
         "TESTING STRATEGY", [
             "Coverage completeness — every AC has P/N/E cases",
             "Duplicate detection — semantic similarity check",
             "Traceability — every TC → Story → AC linked",
             "Golden dataset comparison — >= 85% coverage match",
         ], accent=LIGHT_BLUE)

# Coverage targets table
add_table(slide, Inches(0.5), Inches(3.9), Inches(12.333), Inches(2.2),
          ["Coverage Dimension", "What We Check", "Target"],
          [
              ["AC Coverage", "Every acceptance criterion has at least 1 test case", "100%"],
              ["Positive Path", "Each AC has a happy-path test", "100%"],
              ["Negative Path", "Each AC has at least 1 failure-mode test", ">= 90%"],
              ["Edge Cases", "Boundary values, empty inputs, max lengths", ">= 80%"],
              ["Business Rules", "Each rule has violation scenarios", ">= 85%"],
          ])

add_text_box(slide, Inches(0.5), Inches(6.35), Inches(12), Inches(0.4),
             "Output:  Structured, auditable JSON test cases with confidence scores and Jira traceability",
             font_size=14, color=GREEN_OK)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — MODULE 4: AUTOMATION SCRIPT AGENT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Module 4: Automation Script Agent", "From Test Cases to Executable Code", 8)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.9), Inches(2.5),
         "CODE GENERATION", [
             "UI → Playwright / Selenium",
             "API → REST Assured / Supertest",
             "DB → Parameterized SQL queries",
             "Follows Page Object Model",
             "Uses stable locators (data-testid)",
         ], accent=ACCENT_TEAL)

add_card(slide, Inches(4.7), Inches(1.5), Inches(3.9), Inches(2.5),
         "QUALITY CHECKS", [
             "Syntax valid — compiles clean",
             "Locators robust (score >= 7/10)",
             "Assertions match expected results",
             "No hard-coded waits (sleep)",
             "POM structure compliance",
         ], accent=LIGHT_BLUE)

add_card(slide, Inches(8.9), Inches(1.5), Inches(3.9), Inches(2.5),
         "SELF-HEALING", [
             "Broken locator → find by text/role",
             "App flow changed → regenerate script",
             "Auto-heal confidence > 0.8 → apply",
             "Auto-heal confidence < 0.8 → flag",
             "Recovery: minutes, not hours",
         ], accent=GREEN_OK)

# Locator robustness visual
add_shape_bg(slide, Inches(0.5), Inches(4.4), Inches(12.333), Inches(2.2), RGBColor(0x0D, 0x19, 0x33))
add_text_box(slide, Inches(0.8), Inches(4.55), Inches(5), Inches(0.4),
             "LOCATOR ROBUSTNESS SCALE", font_size=16, color=ACCENT_TEAL, bold=True)

add_table(slide, Inches(0.8), Inches(4.95), Inches(7), Inches(1.5),
          ["Locator Strategy", "Score", "Use When"],
          [
              ["data-testid", "10/10 ✓", "Always preferred"],
              ["id attribute", "8/10 ✓", "If IDs are stable"],
              ["CSS (class-based)", "5/10 ⚠", "Risky — classes change"],
              ["XPath (absolute)", "2/10 ✗", "Never acceptable"],
          ])

add_table(slide, Inches(8.5), Inches(4.95), Inches(4.2), Inches(1.0),
          ["Key Metric", "Target"],
          [
              ["Compilation Success", ">= 95%"],
              ["Auto-Heal Success", ">= 70%"],
          ])


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — MODULE 5: EXECUTION ENGINE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Module 5: Execution Engine", "Running Tests at Scale, Reliably", 9)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.9), Inches(2.2),
         "WHAT IT DOES", [
             "Execute across Chrome, Firefox, Edge",
             "Parallel execution (10x workers)",
             "Environment management (QA/Stage)",
             "CI/CD pipeline integration",
         ], accent=ACCENT_TEAL)

add_card(slide, Inches(4.7), Inches(1.5), Inches(3.9), Inches(2.2),
         "RETRY POLICY", [
             "Element not found → retry 2x",
             "Network timeout → retry 3x + backoff",
             "Browser crash → restart + retry 2x",
             "Assertion failure → NO retry (real bug)",
         ], accent=ORANGE_WARN)

add_card(slide, Inches(8.9), Inches(1.5), Inches(3.9), Inches(2.2),
         "CHAOS TESTING", [
             "Kill browser mid-test → partial save",
             "Network disconnect → retry + alert",
             "Memory pressure → clean shutdown",
             "Grid node removed → redistribute",
         ], accent=RED_ALERT)

# Metrics
add_table(slide, Inches(0.5), Inches(4.1), Inches(12.333), Inches(2.2),
          ["Metric", "Target", "Alert Threshold"],
          [
              ["Execution Success Rate", ">= 90%", "< 85%"],
              ["Flakiness %", "< 5%", "> 10%"],
              ["Retry Recovery Rate", ">= 60%", "< 40%"],
              ["Parallel Efficiency", ">= 70%", "< 50%"],
              ["Avg Execution Time / Test", "< 45 seconds", "> 90 seconds"],
          ])


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — MODULE 6: RESULTS & RCA
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Module 6: Results & RCA Agent", "From Failures to Root Causes to Jira Defects", 10)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.9), Inches(2.0),
         "EVIDENCE CAPTURE", [
             "Console & server logs",
             "Screenshot at failure point",
             "Video recording of full test",
             "Network HAR trace + DOM snapshot",
         ], accent=ACCENT_TEAL)

add_card(slide, Inches(4.7), Inches(1.5), Inches(3.9), Inches(2.0),
         "RCA CLASSIFICATION", [
             "Assertion mismatch → App Bug",
             "Element not found → Locator Issue",
             "HTTP 500 → Backend Bug",
             "Timeout → Infra Issue",
         ], accent=LIGHT_BLUE)

add_card(slide, Inches(8.9), Inches(1.5), Inches(3.9), Inches(2.0),
         "AUTO JIRA DEFECTS", [
             "Bug title + steps to reproduce",
             "Evidence attached (screenshots, logs)",
             "Severity mapped from test priority",
             "Duplicate detection (fingerprinting)",
         ], accent=PURPLE)

# Key metrics + false positive table
add_shape_bg(slide, Inches(0.5), Inches(3.9), Inches(12.333), Inches(2.7), RGBColor(0x0D, 0x19, 0x33))

add_table(slide, Inches(0.8), Inches(4.1), Inches(5.5), Inches(1.5),
          ["Metric", "Target"],
          [
              ["RCA Accuracy %", ">= 80%"],
              ["False Positive Rate", "< 10%"],
              ["Evidence Completeness", ">= 98%"],
              ["Correct Severity Assignment", ">= 90%"],
          ])

add_text_box(slide, Inches(7.0), Inches(4.2), Inches(5.5), Inches(0.4),
             "FALSE POSITIVE FILTERING", font_size=16, color=ORANGE_WARN, bold=True)
add_text_box(slide, Inches(7.0), Inches(4.7), Inches(5.5), Inches(1.5),
             "▸  Infra failure → NOT filed as app bug\n"
             "▸  Flaky test (passes on retry) → NOT filed\n"
             "▸  Environment config error → NOT filed\n"
             "▸  Only real app bugs create Jira defects",
             font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC))


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — AI METRICS FRAMEWORK
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Module 7: AI Metrics Framework", '"AI quality is measured, not assumed"', 11)

# Four metric cards
metrics = [
    ("Requirement\nAccuracy", ">= 85%", "How well AC is extracted\nvs golden dataset", ACCENT_TEAL),
    ("Test Coverage\nScore", ">= 85%", "Positive + Negative +\nEdge case completeness", LIGHT_BLUE),
    ("Hallucination\nRate", "< 5%", "AI-invented content\nwith no source in story", RED_ALERT),
    ("Decision\nConfidence", ">= 0.85", "Agent self-reported\ncertainty per output", ACCENT_BLUE),
]

for i, (title, target, desc, color) in enumerate(metrics):
    x = Inches(0.5) + Inches(i * 3.15)
    add_shape_bg(slide, x, Inches(1.5), Inches(2.95), Inches(2.7), RGBColor(0x1A, 0x2A, 0x44))
    add_shape_bg(slide, x, Inches(1.5), Inches(2.95), Inches(0.05), color)
    add_text_box(slide, x + Inches(0.15), Inches(1.7), Inches(2.65), Inches(0.65),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(2.4), Inches(2.65), Inches(0.5),
                 target, font_size=28, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(3.0), Inches(2.65), Inches(0.8),
                 desc, font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))

# Confidence routing
add_shape_bg(slide, Inches(0.5), Inches(4.6), Inches(12.333), Inches(2.0), RGBColor(0x0D, 0x19, 0x33))
add_text_box(slide, Inches(0.8), Inches(4.75), Inches(5), Inches(0.4),
             "CONFIDENCE-BASED ROUTING", font_size=16, color=ACCENT_TEAL, bold=True)

add_table(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.2),
          ["Confidence Score", "Routing Decision", "Human Involvement"],
          [
              [">= 0.85", "Auto-proceed to next agent", "None"],
              ["0.70 – 0.84", "Proceed with flag for optional review", "Optional"],
              ["< 0.70", "BLOCK — Require human approval", "Mandatory"],
          ])


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — EXECUTION METRICS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Automation Execution Metrics", "Measuring Runtime Performance & Stability", 12)

# Execution metrics
add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.3),
         "EXECUTION METRICS", [
             "Pass / Fail Rate — target >= 90% pass",
             "Retry Recovery Rate — target >= 60%",
             "Avg Execution Time — target < 45s per test",
             "Parallel Efficiency — target >= 70%",
             "Suite Completion Rate — target >= 98%",
         ], accent=ACCENT_TEAL)

# Stability metrics
add_card(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.3),
         "STABILITY METRICS", [
             "Flaky Test Rate — target < 5%",
             "  Root: timing 45%, data 25%, env 20%, order 10%",
             "Auto-Heal Success — target >= 70%",
             "Script Regeneration Success — target >= 80%",
             "Reduce flakiness by 20%/month goal",
         ], accent=GREEN_OK)

# Correlation insight
add_shape_bg(slide, Inches(0.5), Inches(4.2), Inches(12.333), Inches(2.4), RGBColor(0x0D, 0x19, 0x33))
add_text_box(slide, Inches(0.8), Inches(4.35), Inches(5), Inches(0.4),
             "CORRELATION INSIGHTS", font_size=16, color=ORANGE_WARN, bold=True)

add_table(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.6),
          ["AI Quality", "Execution Quality", "Diagnosis", "Action"],
          [
              ["High Accuracy", "Low Pass Rate", "Environment / Infra problem", "Fix infra, not tests"],
              ["Low Accuracy", "High Pass Rate", "Coverage gap — wrong things tested", "Improve AI prompts"],
              ["High Hallucination", "High Pass Rate", "False confidence", "Audit test cases"],
          ])


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — SUPERVISOR AGENT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Supervisor / Orchestrator Agent", '"This agent makes the system truly autonomous"', 13)

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.9), Inches(2.3),
         "COORDINATES", [
             "Agent A finishes → trigger Agent B",
             "Pass outputs between agents",
             "Manage execution dependencies",
             "Handle concurrent pipelines",
         ], accent=ACCENT_TEAL)

add_card(slide, Inches(4.7), Inches(1.5), Inches(3.9), Inches(2.3),
         "DECIDES", [
             ">= 0.85 confidence → Auto-proceed",
             "0.70–0.84 → Proceed with flag",
             "< 0.70 → Escalate to human",
             "Retry vs Skip vs Abort logic",
         ], accent=LIGHT_BLUE)

add_card(slide, Inches(8.9), Inches(1.5), Inches(3.9), Inches(2.3),
         "ESCALATES", [
             "Low confidence on P0 story",
             "3+ consecutive agent failures",
             "Hallucination detected",
             "Auto-created Blocker defect",
         ], accent=RED_ALERT)

# Metrics + Pipeline state
add_shape_bg(slide, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.4), RGBColor(0x0D, 0x19, 0x33))
add_text_box(slide, Inches(0.8), Inches(4.35), Inches(5), Inches(0.4),
             "PIPELINE STATES", font_size=16, color=ACCENT_TEAL, bold=True)
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(5.2), Inches(1.5),
             "INGESTING → INTERPRETING → DESIGNING\n"
             "    → SCRIPTING → EXECUTING → ANALYZING\n"
             "        → REPORTING → COMPLETE\n\n"
             "Each transition: logged, audited, reversible",
             font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC))

add_shape_bg(slide, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.4), RGBColor(0x0D, 0x19, 0x33))
add_text_box(slide, Inches(7.1), Inches(4.35), Inches(5), Inches(0.4),
             "KEY METRICS", font_size=16, color=ACCENT_TEAL, bold=True)
add_table(slide, Inches(7.1), Inches(4.8), Inches(5.2), Inches(1.2),
          ["Metric", "Target"],
          [
              ["Decision Accuracy", ">= 95%"],
              ["Pipeline Completion Rate", ">= 95%"],
              ["Escalation Rate", "< 15%"],
          ])


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — FEEDBACK LOOP
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Feedback Loop & Continuous Learning", "The System Gets Better Over Time", 14)

# The loop visual
add_shape_bg(slide, Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.2), RGBColor(0x0D, 0x19, 0x33))
add_text_box(slide, Inches(0.8), Inches(1.55), Inches(12), Inches(0.4),
             "THE LEARNING LOOP", font_size=18, color=ACCENT_TEAL, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(12), Inches(0.5),
             "Execution Results   →   Metrics Analysis   →   Prompt Tuning   →   Better Output   →   Repeat",
             font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Three areas
add_card(slide, Inches(0.5), Inches(3.0), Inches(3.9), Inches(2.2),
         "WHAT GETS TUNED", [
             "LLM prompts (rules, examples)",
             "Confidence thresholds",
             "Retry strategies",
             "Agent decision boundaries",
         ], accent=ACCENT_TEAL)

add_card(slide, Inches(4.7), Inches(3.0), Inches(3.9), Inches(2.2),
         "TESTING STRATEGY", [
             "Prompt regression testing",
             "  — every change vs golden dataset",
             "Model drift detection",
             "  — daily canary tests",
         ], accent=LIGHT_BLUE)

add_card(slide, Inches(8.9), Inches(3.0), Inches(3.9), Inches(2.2),
         "SAFEGUARDS", [
             "Prompt version control (Git)",
             "Rollback in < 5 minutes",
             "No deploy without regression pass",
             "Monthly benchmark comparisons",
         ], accent=PURPLE)

# Key insight
add_shape_bg(slide, Inches(0.5), Inches(5.6), Inches(12.333), Inches(1.0), RGBColor(0x14, 0x1E, 0x3A))
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(12), Inches(0.7),
             "Key Insight:  Treat prompts like code — version control, test, review, deploy, rollback.\n"
             "A one-word prompt change can shift AI behavior dramatically.",
             font_size=15, color=ORANGE_WARN)


# ══════════════════════════════════════════════════════════════
# SLIDE 15 — DATA FLOW (BONUS)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "End-to-End Data Flow", "From Jira Story to Measurable Test Results", 15)

# Flow steps
flow_steps = [
    ("1", "JIRA\nSTORY", ACCENT_TEAL),
    ("2", "INGEST &\nVALIDATE", ACCENT_TEAL),
    ("3", "UNDERSTAND\nREQUIREMENTS", LIGHT_BLUE),
    ("4", "DESIGN\nTEST CASES", LIGHT_BLUE),
    ("5", "GENERATE\nSCRIPTS", ACCENT_BLUE),
    ("6", "EXECUTE\nTESTS", ACCENT_BLUE),
    ("7", "ANALYZE\n& RCA", PURPLE),
    ("8", "METRICS &\nFEEDBACK", GREEN_OK),
]

for i, (num, label, color) in enumerate(flow_steps):
    x = Inches(0.4) + Inches(i * 1.6)
    y = Inches(1.8)
    # Circle/box
    add_shape_bg(slide, x, y, Inches(1.35), Inches(1.5), RGBColor(0x1A, 0x2A, 0x44))
    add_shape_bg(slide, x, y, Inches(1.35), Inches(0.05), color)
    add_text_box(slide, x, y + Inches(0.15), Inches(1.35), Inches(0.35),
                 num, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.55), Inches(1.25), Inches(0.8),
                 label, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < len(flow_steps) - 1:
        add_text_box(slide, x + Inches(1.35), y + Inches(0.5), Inches(0.3), Inches(0.4),
                     "→", font_size=22, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# What happens at each step
add_shape_bg(slide, Inches(0.5), Inches(3.7), Inches(12.333), Inches(2.8), RGBColor(0x0D, 0x19, 0x33))

details = [
    ("At Each Step:", ACCENT_TEAL, [
        "Input validated against schema",
        "Output carries confidence score",
        "Supervisor monitors and decides",
        "Metrics captured for dashboard",
    ]),
    ("Quality Gates:", GREEN_OK, [
        "Confidence < 0.70 → Human review",
        "Hallucination detected → Pipeline pause",
        "Compilation fails → Auto-retry/regen",
        "RCA flags infra → No false bug filed",
    ]),
    ("Feedback:", PURPLE, [
        "Every result feeds back to prompts",
        "Model drift detected by daily canaries",
        "Monthly benchmarks track improvement",
        "Golden datasets provide ground truth",
    ]),
]

for i, (title, color, items) in enumerate(details):
    x = Inches(0.8) + Inches(i * 4.1)
    add_text_box(slide, x, Inches(3.85), Inches(3.8), Inches(0.4),
                 title, font_size=15, color=color, bold=True)
    y = Inches(4.3)
    for item in items:
        add_text_box(slide, x, y, Inches(3.8), Inches(0.3),
                     "▸  " + item, font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))
        y += Inches(0.3)


# ══════════════════════════════════════════════════════════════
# SLIDE 16 — COMPLETE METRICS DASHBOARD
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
slide_header(slide, "Complete Metrics Dashboard", "AI Quality + Execution Quality + Operations", 16)

# AI metrics
add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(0.4),
             "AI QUALITY METRICS", font_size=16, color=ACCENT_TEAL, bold=True)
add_table(slide, Inches(0.5), Inches(1.95), Inches(5.8), Inches(1.6),
          ["Metric", "Target", "Status"],
          [
              ["Requirement Accuracy", ">= 85%", "87.3% ✓"],
              ["Test Coverage Score", ">= 85%", "83.6% ⚠"],
              ["Hallucination Rate", "< 5%", "3.2% ✓"],
              ["Decision Confidence", ">= 0.85", "0.88 ✓"],
          ])

# Execution metrics
add_text_box(slide, Inches(6.8), Inches(1.5), Inches(4), Inches(0.4),
             "EXECUTION METRICS", font_size=16, color=LIGHT_BLUE, bold=True)
add_table(slide, Inches(6.8), Inches(1.95), Inches(5.8), Inches(1.6),
          ["Metric", "Target", "Status"],
          [
              ["Pass Rate", ">= 90%", "88.4% ⚠"],
              ["Flakiness", "< 5%", "4.6% ✓"],
              ["Auto-Heal Success", ">= 70%", "73% ✓"],
              ["Parallel Efficiency", ">= 70%", "83% ✓"],
          ])

# Operations metrics
add_text_box(slide, Inches(0.5), Inches(3.9), Inches(4), Inches(0.4),
             "OPERATIONS METRICS", font_size=16, color=PURPLE, bold=True)
add_table(slide, Inches(0.5), Inches(4.35), Inches(5.8), Inches(1.2),
          ["Metric", "Target", "Status"],
          [
              ["Pipeline Completion", ">= 95%", "96.2% ✓"],
              ["Mean Time per Story", "< 15 min", "12 min ✓"],
              ["Escalation Rate", "< 15%", "12% ✓"],
          ])

# Trend
add_text_box(slide, Inches(6.8), Inches(3.9), Inches(4), Inches(0.4),
             "LEARNING TREND (3 MONTHS)", font_size=16, color=GREEN_OK, bold=True)
add_table(slide, Inches(6.8), Inches(4.35), Inches(5.8), Inches(1.2),
          ["Metric", "Month 1", "Month 3", "Trend"],
          [
              ["Req Accuracy", "78%", "87%", "▲ +9%"],
              ["Hallucination", "9%", "3%", "▼ -6% ✓"],
              ["Pass Rate", "80%", "89%", "▲ +9%"],
          ])

# Bottom insight
add_shape_bg(slide, Inches(0.5), Inches(5.9), Inches(12.333), Inches(0.7), RGBColor(0x14, 0x1E, 0x3A))
add_text_box(slide, Inches(0.8), Inches(5.95), Inches(12), Inches(0.5),
             "Every metric has: a target, an alert threshold, an associated action, and a trend line.",
             font_size=15, color=ORANGE_WARN, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 17 — CLOSING SLIDE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SECTION_BG)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RGBColor(0x00, 0x2B, 0x5C), alpha=40)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_TEAL)

# Benefits
add_text_box(slide, Inches(1.0), Inches(0.6), Inches(11), Inches(0.6),
             "Architecture Value Proposition", font_size=34, color=WHITE, bold=True)
add_shape_bg(slide, Inches(1.0), Inches(1.2), Inches(3), Inches(0.04), ACCENT_TEAL)

benefits = [
    ("SCALABLE", "10 or 10,000 stories\nSame platform, no extra headcount", GREEN_OK),
    ("SELF-HEALING", "Broken locators auto-fixed\nFailed tests auto-retried", LIGHT_BLUE),
    ("TRUSTWORTHY AI", "Every decision measured\nAudited and explainable", ACCENT_TEAL),
    ("PRODUCTION-READY", "CI/CD integrated\nSecurity hardened, monitoring live", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(benefits):
    x = Inches(0.5) + Inches(i * 3.15)
    y = Inches(1.6)
    add_shape_bg(slide, x, y, Inches(2.95), Inches(1.8), RGBColor(0x14, 0x1E, 0x3A))
    add_shape_bg(slide, x, y, Inches(2.95), Inches(0.05), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.65), Inches(0.4),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.65), Inches(0.8),
                 desc, font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

# ROI
add_shape_bg(slide, Inches(0.5), Inches(3.7), Inches(12.333), Inches(1.3), RGBColor(0x14, 0x1E, 0x3A))
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(4), Inches(0.4),
             "ROI SUMMARY", font_size=18, color=GREEN_OK, bold=True)

roi_items = [
    ("Test design time:", "Days → Minutes", "95% reduction"),
    ("Script maintenance:", "Constant → Minimal", "70% reduction"),
    ("Bug escape rate:", "15% → < 5%", "67% reduction"),
]
y = Inches(4.25)
for label, change, impact in roi_items:
    add_text_box(slide, Inches(0.8), y, Inches(2.5), Inches(0.3),
                 label, font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB), bold=True)
    add_text_box(slide, Inches(3.3), y, Inches(3.5), Inches(0.3),
                 change, font_size=14, color=WHITE)
    add_text_box(slide, Inches(7.0), y, Inches(3), Inches(0.3),
                 impact, font_size=14, color=GREEN_OK, bold=True)
    y += Inches(0.3)

# Closing statement box
add_shape_bg(slide, Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.3), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.4), Inches(11.8), Inches(1.0),
             '"This architecture ensures confidence in both the application\n'
             'under test and the AI testing platform itself."',
             font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 17)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = r"F:\work\testing-tool\Agentic-AI-Testing-Architecture.pptx"
prs.save(output_path)
print(f"[OK] Presentation saved: {output_path}")
print(f"     17 slides created")
